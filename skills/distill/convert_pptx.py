#!/usr/bin/env python3
"""Convert a PowerPoint file to slide-structured Markdown.

Usage: convert_pptx.py --input <path.pptx> --output <path.md>
"""

import argparse
import sys
from pathlib import Path

try:
    from pptx import Presentation
except ImportError:
    print("Error: python-pptx is not installed.", file=sys.stderr)
    print("Install with: pip install python-pptx==1.0.2", file=sys.stderr)
    sys.exit(1)


def extract_text_from_shape(shape):
    """Extract text from a shape, handling text frames and tables."""
    lines = []

    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            text = paragraph.text.strip()
            if text:
                lines.append(text)

    if shape.has_table:
        table = shape.table
        rows = []
        for row in table.rows:
            cells = []
            for cell in row.cells:
                cells.append(cell.text.strip())
            rows.append(cells)

        if rows:
            # Build markdown table
            header = "| " + " | ".join(rows[0]) + " |"
            separator = "| " + " | ".join("---" for _ in rows[0]) + " |"
            lines.append(header)
            lines.append(separator)
            for row in rows[1:]:
                lines.append("| " + " | ".join(row) + " |")

    return "\n".join(lines)


def convert_pptx_to_markdown(input_path: Path) -> str:
    """Convert a PPTX file to slide-structured Markdown."""
    prs = Presentation(str(input_path))
    slides_md = []

    for i, slide in enumerate(prs.slides, 1):
        # Extract title
        title = ""
        if slide.shapes.title:
            title = slide.shapes.title.text.strip()

        # Extract body content (skip the title shape)
        body_parts = []
        for shape in slide.shapes:
            if shape == slide.shapes.title:
                continue
            text = extract_text_from_shape(shape)
            if text:
                body_parts.append(text)

        # Extract speaker notes
        notes = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                notes = notes_text

        # Build slide markdown
        heading = f"## Slide {i}: {title}" if title else f"## Slide {i}"
        parts = [heading]

        if body_parts:
            parts.append("\n".join(body_parts))

        if notes:
            parts.append(f"\n> Speaker notes: {notes}")

        slides_md.append("\n\n".join(parts))

    return "\n\n---\n\n".join(slides_md) + "\n"


def main():
    parser = argparse.ArgumentParser(description="Convert PPTX to Markdown")
    parser.add_argument("--input", required=True, help="Path to input .pptx file")
    parser.add_argument("--output", required=True, help="Path to output .md file")
    args = parser.parse_args()

    input_path = Path(args.input)
    output_path = Path(args.output)

    if not input_path.exists():
        print(f"Error: Input file not found: {input_path}", file=sys.stderr)
        sys.exit(1)

    if not input_path.suffix.lower() == ".pptx":
        print(f"Error: Expected .pptx file, got: {input_path.suffix}", file=sys.stderr)
        sys.exit(1)

    try:
        markdown = convert_pptx_to_markdown(input_path)
        output_path.write_text(markdown, encoding="utf-8")
        print(f"Converted: {input_path} -> {output_path}")
    except Exception as e:
        print(f"Error converting {input_path}: {e}", file=sys.stderr)
        sys.exit(1)


if __name__ == "__main__":
    main()
